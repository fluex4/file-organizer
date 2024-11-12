import os
import shutil
import PyPDF2
import fitz
import docx
import sys
from pptx import Presentation
import joblib
import pandas as pd
import nltk
from nltk.tokenize import word_tokenize
from nltk.util import ngrams
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.naive_bayes import MultinomialNB
from sklearn.model_selection import train_test_split
from sklearn.metrics import confusion_matrix, classification_report, accuracy_score
from wordcloud import WordCloud
import matplotlib.pyplot as plt

class TextFileClassifier:
    BASE_DIR = os.path.join(os.path.expanduser('~'), "File Organizer")
    MODEL_DIR = os.path.join(BASE_DIR, "models")
    DATA_DIR = os.path.join(BASE_DIR, "data")

    def __init__(self, log_filename="classification_log.txt"):
        os.makedirs(self.MODEL_DIR, exist_ok=True)
        os.makedirs(self.DATA_DIR, exist_ok=True)

        self.vectorizer_path = os.path.join(self.MODEL_DIR, "vectorizer.pkl")
        self.model_path = os.path.join(self.MODEL_DIR, "classifier_model.pkl")
        self.log_path = os.path.join(os.getcwd(), log_filename)

        if hasattr(sys, '_MEIPASS'):
            self.source_data_path = os.path.join(sys._MEIPASS, 'data.csv')
        else:
            self.source_data_path = "data.csv"

        nltk.download('stopwords')
        nltk.download('punkt')

        if os.path.exists(self.model_path) and os.path.exists(self.vectorizer_path):
            self.vectorizer = joblib.load(self.vectorizer_path)
            self.clf = joblib.load(self.model_path)
            print("Loaded model and vectorizer from disk.")
        else:
            self.vectorizer = None
            self.clf = None
            print("Model and vectorizer not found. Please train the classifier.")

    def save_model(self):
        joblib.dump(self.vectorizer, self.vectorizer_path)
        joblib.dump(self.clf, self.model_path)
        print("Model and vectorizer saved to disk.")

    def read_text_from_file(self, file_path):
        if file_path.endswith('.pdf'):
            return self.extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            return self.read_word(file_path)
        elif file_path.endswith('.pptx'):
            return self.read_ppt(file_path)
        elif file_path.endswith('.txt'):
            return self.read_txt(file_path)
        elif file_path.endswith('.py'):
            return self.read_txt(file_path)
        elif file_path.endswith('.js'):
            return self.read_txt(file_path)
        elif file_path.endswith('.c'):
            return self.read_txt(file_path)
        else:
            print(f"Unsupported file type for {file_path}")
            return None

    def read_pdf(self, file_path):
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text
    
    def extract_text_from_pdf(self,file_path):
        text = ""
        pdf_document = fitz.open(file_path)
        for page_num in range(min(6,pdf_document.page_count)):
            page = pdf_document[page_num]
            text += page.get_text()
        return text

    def read_word(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def read_ppt(self, file_path):
        text = ""
        presentation = Presentation(file_path)
        for slide in min(6,presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def read_txt(self, file_path):
        with open(file_path, 'r') as file:
            return file.read()

    def classify_files_in_folder(self, folder_path):
        # folder_path=os.path.abspath(folder_path)
        folder_path=os.getcwd()
        print(folder_path)
        if not os.path.exists(folder_path):
            print(f"Folder {folder_path} does not exist.")
            return

        for filename in os.listdir(folder_path):
            try:
                file_path = os.path.join(folder_path, filename)
                file_content = self.read_text_from_file(file_path)
                if file_content:
                    predicted_category = self.predict_category(file_content)
                    category_folder = os.path.join(folder_path, predicted_category)
                    if not os.path.exists(category_folder):
                        os.makedirs(category_folder)
                    destination_path = os.path.join(category_folder, filename)
                    shutil.move(file_path, destination_path)
                    self.log_movement(file_path, destination_path)
                else:
                    print(f"Skipping unsupported file type for {filename}")
            except:
                continue

    def log_movement(self, original_path, new_path):
        with open(self.log_path, "a") as log_file:
            log_file.write(f"'{original_path}' - '{new_path}'\n")

    def process_text(self, text):
        words = word_tokenize(text.lower())
        stop_words = set(nltk.corpus.stopwords.words('english'))
        filtered_words = [word for word in words if word.isalnum() and word not in stop_words]
        return filtered_words

    def analyze_frequencies(self, words):
        freq_dist = nltk.FreqDist(words)
        bigrams = ngrams(words, 2)
        bigram_freq = nltk.FreqDist(bigrams)
        trigrams = ngrams(words, 3)
        trigram_freq = nltk.FreqDist(trigrams)
        return freq_dist, bigram_freq, trigram_freq

    def visualize_frequencies(self, freq_dist):
        wordcloud = WordCloud(width=800, height=400).generate_from_frequencies(freq_dist)
        plt.figure(figsize=(10, 5))
        plt.imshow(wordcloud, interpolation='bilinear')
        plt.axis('off')
        plt.show()

    def train_classifier(self, texts, labels):
        self.vectorizer = TfidfVectorizer(ngram_range=(1, 3))
        X = self.vectorizer.fit_transform(texts)
        X_train, X_test, y_train, y_test = train_test_split(X, labels, test_size=0.2, random_state=40)
        self.clf = MultinomialNB()
        self.clf.fit(X_train, y_train)
        y_pred = self.clf.predict(X_test)
        print("\nClassification Report:\n", classification_report(y_test, y_pred))
        print("\nConfusion Matrix:\n", confusion_matrix(y_test, y_pred))
        print("\nAccuracy:\n", accuracy_score(y_test, y_pred))
        self.save_model()

    def predict_category(self, text):
        if self.vectorizer is None or self.clf is None:
            print("Model and vectorizer are not loaded. Please train the classifier first.")
            return None
        text_vector = self.vectorizer.transform([text])
        return self.clf.predict(text_vector)[0]
